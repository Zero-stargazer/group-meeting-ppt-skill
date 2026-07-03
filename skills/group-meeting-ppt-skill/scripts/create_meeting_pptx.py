#!/usr/bin/env python
"""Create a passable Chinese group-meeting PPTX from structured paper notes."""

from __future__ import annotations

import argparse
import sys
import textwrap
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - exercised by CLI users
    raise SystemExit(
        "Missing dependency: python-pptx. "
        "Use the Codex bundled Python or install it with: python -m pip install python-pptx"
    ) from exc

SCRIPT_DIR = Path(__file__).resolve().parent
if str(SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPT_DIR))

from create_meeting_pack import as_list, read_json, slide_plan, text  # noqa: E402


NAVY = RGBColor(18, 28, 45)
BLUE = RGBColor(44, 103, 255)
CYAN = RGBColor(33, 203, 255)
MINT = RGBColor(0, 173, 181)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(244, 247, 251)
MID = RGBColor(106, 118, 140)
DARK = RGBColor(37, 45, 60)
PALE_BLUE = RGBColor(229, 240, 255)


def cm(value: float) -> float:
    return value / 2.54


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text_value: str,
    *,
    font_size: int = 20,
    bold: bool = False,
    color: RGBColor = DARK,
    align=PP_ALIGN.LEFT,
    font_name: str = "Microsoft YaHei",
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text_value
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_multiline(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: list[str],
    *,
    font_size: int = 18,
    color: RGBColor = DARK,
    bullet: bool = False,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        if bullet:
            paragraph.text = f"• {line}"
    return box


def add_header(slide, slide_id: str, title: str, total: int = 10) -> None:
    add_textbox(slide, 0.55, 0.32, 1.15, 0.33, slide_id, font_size=12, bold=True, color=BLUE)
    add_textbox(slide, 1.62, 0.22, 8.8, 0.55, title, font_size=25, bold=True, color=NAVY)
    add_textbox(slide, 11.15, 0.31, 1.45, 0.35, f"{slide_id[-2:]}/{total}", font_size=12, color=MID, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.86), Inches(12.25), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = PALE_BLUE
    line.line.color.rgb = PALE_BLUE


def add_label(slide, left: float, top: float, label: str, color: RGBColor = BLUE) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(1.25), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_panel(slide, left: float, top: float, width: float, height: float, *, fill=LIGHT, line=PALE_BLUE) -> None:
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    panel.line.color.rgb = line


def short_lines(value: str, width: int = 28, max_lines: int = 4) -> list[str]:
    value = " ".join(str(value).split())
    if not value:
        return ["待根据原文补充。"]
    wrapped = textwrap.wrap(value, width=width, break_long_words=True, replace_whitespace=False)
    if len(wrapped) > max_lines:
        wrapped = wrapped[: max_lines - 1] + [wrapped[max_lines - 1] + "……"]
    return wrapped


def figure_cards(data: dict[str, Any]) -> list[dict[str, str]]:
    cards: list[dict[str, str]] = []
    for fig in as_list(data, "figures_to_discuss"):
        if isinstance(fig, dict):
            cards.append(
                {
                    "label": str(fig.get("label", "图表")).strip(),
                    "why": str(fig.get("why_use", "说明论文关键证据。")).strip(),
                    "say": str(fig.get("what_to_say", "解释这张图支持了哪个结论。")).strip(),
                }
            )
    return cards


def visual_suggestion(data: dict[str, Any], index: int) -> tuple[str, list[str]]:
    figures = figure_cards(data)
    if index in (5, 6) and figures:
        fig = figures[min(index - 5, len(figures) - 1)]
        return fig["label"], short_lines(f'{fig["why"]}；{fig["say"]}', width=19, max_lines=5)
    if index == 3:
        route = as_list(data, "method_route")
        return "方法路线", [f"{i + 1}. {item}" for i, item in enumerate(route[:4])] or ["3—5步方法路线"]
    if index == 4:
        return "流程图占位", ["输入 → 处理 → 读出 → 判定", "优先放论文 workflow / schematic"]
    if index == 7:
        return "结果表", [str(item) for item in as_list(data, "key_findings")[:3]] or ["列出 2—3 个关键结果"]
    if index == 8:
        return "风险清单", [str(item) for item in as_list(data, "limitations")[:3]] or ["样本量 / 对照 / 泛化 / 统计"]
    if index == 9:
        return "迁移价值", ["我能复用什么？", "第一步实验怎么设计？"]
    if index == 10:
        return "追问预案", ["可靠性", "样本量", "对照实验", "泛化能力"]
    return "核心信息", ["1句话结论", "2—3条证据", "1个需要核对的点"]


def build_deck(data: dict[str, Any], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    title_zh = text(data, "paper_title_zh", "论文题目待补充")
    title_en = text(data, "paper_title_en", "")
    direction = text(data, "research_direction", "研究方向待补充")
    meeting_type = text(data, "meeting_type", "文献组会")
    duration = text(data, "duration_minutes", "15")

    for idx, (slide_id, slide_title, goal, notes) in enumerate(slide_plan(data), start=1):
        slide = prs.slides.add_slide(blank)
        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = WHITE
        add_header(slide, slide_id, slide_title)

        if idx == 1:
            add_textbox(slide, 0.75, 1.25, 8.1, 0.45, meeting_type, font_size=18, bold=True, color=BLUE)
            add_textbox(slide, 0.75, 1.82, 9.25, 1.15, title_zh, font_size=31, bold=True, color=NAVY)
            if title_en:
                add_textbox(slide, 0.78, 3.05, 8.95, 0.5, title_en, font_size=14, color=MID)
            add_panel(slide, 0.75, 4.15, 5.6, 1.35)
            add_label(slide, 1.05, 4.38, "一句话")
            add_multiline(slide, 2.55, 4.32, 3.35, 0.85, short_lines(notes, width=18, max_lines=3), font_size=17)
            add_panel(slide, 6.65, 4.15, 5.55, 1.35)
            add_label(slide, 6.95, 4.38, "场景")
            add_multiline(slide, 8.4, 4.34, 3.3, 0.8, [f"{duration}分钟汇报", direction], font_size=17)
        else:
            add_label(slide, 0.72, 1.18, "这一页目的", color=MINT)
            add_multiline(slide, 0.75, 1.63, 5.25, 0.95, short_lines(goal, width=24, max_lines=3), font_size=20)

            add_label(slide, 0.72, 2.75, "你应该讲", color=BLUE)
            add_panel(slide, 0.72, 3.23, 5.9, 2.55)
            add_multiline(slide, 1.05, 3.53, 5.15, 1.95, short_lines(notes, width=27, max_lines=5), font_size=20)

            card_title, card_lines = visual_suggestion(data, idx)
            add_panel(slide, 7.05, 1.25, 5.5, 4.55, fill=RGBColor(248, 251, 255))
            add_textbox(slide, 7.4, 1.58, 4.65, 0.35, card_title, font_size=18, bold=True, color=BLUE)
            placeholder = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(7.45),
                Inches(2.15),
                Inches(4.72),
                Inches(2.05),
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = PALE_BLUE
            placeholder.line.color.rgb = RGBColor(188, 212, 255)
            tf = placeholder.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = "放论文图 / 流程 / 表格"
            r.font.name = "Microsoft YaHei"
            r.font.size = Pt(18)
            r.font.bold = True
            r.font.color.rgb = BLUE
            add_multiline(slide, 7.45, 4.47, 4.72, 0.95, card_lines, font_size=15)

        add_textbox(
            slide,
            0.75,
            6.72,
            11.75,
            0.28,
            "边界提醒：所有数据、图注和结论必须回原文核对；这只是组会草稿，不是最终版。",
            font_size=10,
            color=MID,
        )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", required=True, help="Structured paper-note JSON file.")
    parser.add_argument("--out", required=True, help="Output .pptx path.")
    args = parser.parse_args(argv)

    data = read_json(Path(args.input))
    out_path = Path(args.out)
    build_deck(data, out_path)
    print(f"Wrote PPTX to {out_path.resolve()}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
