import importlib.util
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
SCRIPT = ROOT / "skills" / "group-meeting-ppt-skill" / "scripts" / "create_meeting_pptx.py"
SAMPLE = ROOT / "examples" / "sample-paper-input.json"


class CreateMeetingPptxTest(unittest.TestCase):
    def test_create_meeting_pptx_generates_ten_slide_deck(self):
        if importlib.util.find_spec("pptx") is None:
            self.skipTest("python-pptx is not installed in this Python environment")

        with tempfile.TemporaryDirectory() as tmp:
            out = Path(tmp) / "draft_group_meeting.pptx"
            result = subprocess.run(
                [sys.executable, str(SCRIPT), "--input", str(SAMPLE), "--out", str(out)],
                text=True,
                capture_output=True,
                check=False,
            )

            self.assertEqual(result.returncode, 0, result.stderr)
            self.assertTrue(out.exists())
            self.assertGreater(out.stat().st_size, 10_000)

            from pptx import Presentation

            prs = Presentation(out)
            self.assertEqual(len(prs.slides), 10)
            all_text = "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
            self.assertIn("论文一句话", all_text)
            self.assertIn("老师可能追问", all_text)
            self.assertIn("边界提醒", all_text)


if __name__ == "__main__":
    unittest.main()
